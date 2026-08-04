"""Build the capstone presentation as a real PowerPoint file.

    python scripts/build_deck.py [out.pptx]

Design notes, since the first version of this deck read as machine-written and the
reasons were structural rather than cosmetic:

* Every slide had the same shape -- eyebrow, headline, three paragraphs. Uniform
  structure is the tell. Here the layouts deliberately differ: statement slides
  carry one sentence, the incident slide is full-bleed inverted, two slides are
  native charts, one is a drawn diagram.
* Numbers were set as tables. A table is what you write when you have not decided
  what the number means. The two results slides are charts, so the gap is a thing
  you see rather than a thing you compute.
* Text was everywhere. Slides here are anchors for speech, not the speech.

Palette is a single warm accent on white. White because it projects in a lit
classroom; one accent because a second colour would have to earn itself.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- design tokens

INK = RGBColor(0x14, 0x18, 0x1D)
MUTED = RGBColor(0x66, 0x70, 0x81)
FAINT = RGBColor(0xE3, 0xE7, 0xEB)
ACCENT = RGBColor(0xC2, 0x41, 0x0C)
ACCENT_L = RGBColor(0xFD, 0xBA, 0x74)
GOOD = RGBColor(0x15, 0x80, 0x3D)
BAD = RGBColor(0xB9, 0x1C, 0x1C)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x94, 0xA3, 0xB8)

FONT = "Helvetica Neue"
MONO = "Menlo"

W, H = Inches(13.333), Inches(7.5)
ML = Inches(0.95)  # left margin
MT = Inches(0.85)  # top margin
CW = W - ML - Inches(0.95)  # content width


def deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, colour):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = colour
    r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def text(
    slide,
    x,
    y,
    w,
    h,
    runs,
    *,
    size=18,
    colour=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    spacing=1.15,
    font=FONT,
    anchor=MSO_ANCHOR.TOP,
):
    """runs: a string, or a list of (text, {overrides}) tuples."""
    # Clamp the frame to the canvas. Callers pass a generous height so wrapping has
    # room, which pushed several boxes past the bottom edge; PowerPoint does not clip,
    # so that reads as text hanging off the slide.
    h = min(h, H - y - Inches(0.12))
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    if isinstance(runs, str):
        runs = [(runs, {})]

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    for t, over in runs:
        if t == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = spacing
            continue
        r = p.add_run()
        r.text = t
        f = r.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.italic = over.get("italic", italic)
        f.color.rgb = over.get("colour", colour)
    return box


def eyebrow(slide, label):
    text(
        slide, ML, MT, CW, Inches(0.3), label.upper(), size=11, colour=ACCENT, bold=True
    )


def rule(slide, y, x=None, w=Inches(1.1), colour=ACCENT, thick=Pt(3)):
    x = ML if x is None else x
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thick)
    ln.fill.solid()
    ln.fill.fore_color.rgb = colour
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def headline(slide, y, runs, size=40, w=None):
    return text(
        slide,
        ML,
        y,
        w or Inches(10.2),
        Inches(2.0),
        runs,
        size=size,
        bold=True,
        spacing=1.02,
    )


def body(slide, y, runs, size=17, w=None, colour=MUTED):
    return text(
        slide,
        ML,
        y,
        w or Inches(9.4),
        Inches(2.4),
        runs,
        size=size,
        colour=colour,
        spacing=1.32,
    )


def footer(slide, n, total):
    text(
        slide,
        W - Inches(1.5),
        H - Inches(0.62),
        Inches(0.9),
        Inches(0.3),
        f"{n}",
        size=10,
        colour=SLATE,
        align=PP_ALIGN.RIGHT,
    )


def table(slide, x, y, w, rows, widths, *, header=True, size=13, row_h=Inches(0.42)):
    """Hand-drawn table: python-pptx's native table styling fights the design."""
    cy = y
    for i, row in enumerate(rows):
        cx = x
        is_head = header and i == 0
        for j, cell in enumerate(row):
            content, over = (cell, {}) if isinstance(cell, str) else cell
            text(
                slide,
                cx,
                cy + Inches(0.07),
                widths[j],
                row_h,
                content,
                size=over.get("size", 10 if is_head else size),
                colour=over.get("colour", MUTED if is_head else INK),
                bold=over.get("bold", is_head),
                font=over.get("font", FONT),
            )
            cx += widths[j]
        cy += row_h
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, cy, w, Pt(0.75))
        line.fill.solid()
        line.fill.fore_color.rgb = FAINT if not is_head else SLATE
        line.line.fill.background()
        line.shadow.inherit = False
    return cy


def bar_chart(slide, x, y, w, h, cats, series, colours, *, maximum=None):
    data = CategoryChartData()
    data.categories = cats
    for name, vals in series:
        data.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data)
    ch = gf.chart
    ch.has_legend = len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(12)
        ch.legend.font.name = FONT
    plot = ch.plots[0]
    plot.gap_width = 60
    plot.overlap = -10 if len(series) > 1 else 0

    for si, s in enumerate(ch.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = colours[si % len(colours)]
        s.format.line.fill.background()

    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = FAINT
    va.major_gridlines.format.line.width = Pt(0.75)
    va.tick_labels.font.size = Pt(11)
    va.tick_labels.font.color.rgb = MUTED
    va.tick_labels.font.name = FONT
    if maximum:
        va.maximum_scale = maximum
    va.format.line.fill.background()

    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(12)
    ca.tick_labels.font.color.rgb = INK
    ca.tick_labels.font.name = FONT
    ca.format.line.color.rgb = SLATE
    return ch


# ---------------------------------------------------------------- slides


def build(path: Path) -> Path:
    prs = deck()
    n = 0

    def new(label=None, paper=PAPER):
        nonlocal n
        n += 1
        s = blank(prs)
        bg(s, paper)
        if label:
            eyebrow(s, label)
        return s

    # 1 — title -----------------------------------------------------------
    s = new()
    rule(s, Inches(2.32), w=Inches(1.6), thick=Pt(5))
    text(
        s,
        ML,
        Inches(1.05),
        Inches(11),
        Inches(0.4),
        "CECS 499  ·  SENIOR CAPSTONE  ·  SUMMER 2026",
        size=12,
        colour=ACCENT,
        bold=True,
    )
    headline(
        s,
        Inches(2.65),
        [
            ("Predicting NBA games\n", {}),
            ("without seeing the future", {"colour": ACCENT}),
        ],
        size=54,
        w=Inches(11.5),
    )
    body(
        s,
        Inches(4.85),
        "An agent, a model, and an experiment to find out whether explaining a "
        "prediction makes it better.",
        size=19,
        w=Inches(8.6),
    )
    text(
        s,
        ML,
        Inches(6.15),
        Inches(9),
        Inches(0.8),
        [
            ("Josh Cannon", {"bold": True, "colour": INK}),
            ("  ·  Patrick Haley  ·  Sarvesh Vinod Kumar  ·  Kirtan Patel\n", {}),
            ("Advisor: Prof. Amir Sadovnik", {}),
        ],
        size=14,
        colour=MUTED,
        spacing=1.5,
    )

    # 2 — statement -------------------------------------------------------
    s = new()
    headline(
        s,
        Inches(2.1),
        [
            ("Anyone can predict games badly.\n", {"colour": MUTED}),
            ("Knowing you did it ", {}),
            ("well", {"colour": ACCENT, "italic": True}),
            (" is the hard part.", {}),
        ],
        size=42,
        w=Inches(11.6),
    )
    rule(s, Inches(4.5), w=Inches(1.1))
    body(
        s,
        Inches(4.95),
        [
            (
                "A language model has read the internet. Ask it who won a game in "
                "December 2025 and it may simply ",
                {},
            ),
            ("know", {"bold": True, "colour": INK}),
            (
                ".  Any accuracy number produced that way measures memory, not "
                "prediction — and it looks excellent, which is what makes it dangerous.",
                {},
            ),
        ],
        size=18,
        w=Inches(10.4),
    )

    # 3 — three leaks -----------------------------------------------------
    s = new("Framing")
    headline(s, Inches(1.35), "Three different leaks", size=38)
    body(
        s,
        Inches(2.15),
        "They collapse into one word. They need different defences.",
        size=16,
    )

    cards = [
        ("DATA", "A tool returns an injury\nfiled after the game", "Date gating"),
        (
            "MODEL KNOWLEDGE",
            "The LLM remembers the\nresult from training",
            "Cutoff-pinned model",
        ),
        (
            "BENCHMARK",
            "The system reads the answer\noff what grades it",
            "Remove the tool",
        ),
    ]
    cw, gap = Inches(3.6), Inches(0.42)
    for i, (t, d, fix) in enumerate(cards):
        cx = ML + i * (cw + gap)
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(2.95), cw, Inches(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xF9)
        box.line.color.rgb = FAINT
        box.line.width = Pt(1)
        box.shadow.inherit = False
        box.adjustments[0] = 0.04
        rule(s, Inches(3.25), x=cx + Inches(0.32), w=Inches(0.5), thick=Pt(3))
        text(
            s,
            cx + Inches(0.32),
            Inches(3.55),
            cw - Inches(0.64),
            Inches(0.3),
            t,
            size=10,
            colour=ACCENT,
            bold=True,
        )
        text(
            s,
            cx + Inches(0.32),
            Inches(3.95),
            cw - Inches(0.64),
            Inches(0.9),
            d,
            size=14,
            colour=INK,
            spacing=1.25,
        )
        text(
            s,
            cx + Inches(0.32),
            Inches(4.82),
            cw - Inches(0.64),
            Inches(0.35),
            fix,
            size=13,
            colour=GOOD,
            bold=True,
        )

    body(
        s,
        Inches(5.85),
        [
            ("We anticipated the first two. ", {}),
            ("The third caught us in the act.", {"bold": True, "colour": INK}),
        ],
        size=17,
    )

    # 4 — pipeline diagram ------------------------------------------------
    s = new("Architecture")
    headline(
        s,
        Inches(1.3),
        [("Every query carries an ", {}), ("as-of date", {"colour": ACCENT})],
        size=36,
    )

    steps = [
        (
            "1",
            "gate_snapshot.py  →  data/snapshots/D/",
            "Plain Python. No AI. The future is not filtered — it is absent from disk.",
        ),
        (
            "2",
            "sources.py — every read filtered to ≤ D",
            "Per-tool precision layered on top of the snapshot.",
        ),
        (
            "3",
            "7 tools → agent      8 features → model",
            "Both read through the same gated accessors, so the comparison is fair.",
        ),
        (
            "4",
            "three_arms.py scores all three",
            "Reads results only to score, after the prediction is made.",
        ),
    ]
    y = Inches(2.3)
    for num, title, sub in steps:
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL, ML, y + Inches(0.06), Inches(0.34), Inches(0.34)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        dot.shadow.inherit = False
        text(
            s,
            ML,
            y + Inches(0.11),
            Inches(0.34),
            Inches(0.3),
            num,
            size=12,
            colour=PAPER,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text(
            s,
            ML + Inches(0.62),
            y,
            Inches(10.4),
            Inches(0.32),
            title,
            size=17,
            colour=INK,
            bold=True,
            font=MONO,
        )
        text(
            s,
            ML + Inches(0.62),
            y + Inches(0.38),
            Inches(10.4),
            Inches(0.3),
            sub,
            size=13,
            colour=MUTED,
        )
        y += Inches(1.03)
        if num != "4":
            arr = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                ML + Inches(0.155),
                y - Inches(0.28),
                Pt(1.5),
                Inches(0.22),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT_L
            arr.line.fill.background()
            arr.shadow.inherit = False

    # 5 — why two gates ---------------------------------------------------
    s = new("Why two gates")
    headline(s, Inches(1.3), "They are not redundant", size=38)
    body(
        s,
        Inches(2.2),
        [
            ("A snapshot can only be as strict as its ", {}),
            ("loosest legitimate reader", {"bold": True, "colour": INK}),
            (".", {}),
        ],
        size=19,
        w=Inches(9.5),
    )

    for i, (t, d) in enumerate(
        [
            ("rolling form", "needs games strictly BEFORE D"),
            ("schedule context", "needs games THROUGH D"),
        ]
    ):
        cx = ML + i * Inches(5.0)
        rule(s, Inches(3.15), x=cx, w=Inches(0.7), thick=Pt(3))
        text(
            s,
            cx,
            Inches(3.42),
            Inches(4.4),
            Inches(0.35),
            t,
            size=19,
            colour=INK,
            bold=True,
            font=MONO,
        )
        text(s, cx, Inches(3.92), Inches(4.4), Inches(0.4), d, size=14, colour=MUTED)

    body(
        s,
        Inches(4.75),
        [
            (
                "One on-disk cut cannot serve both without starving one. The snapshot "
                "removes what ",
                {},
            ),
            ("nobody", {"italic": True, "colour": INK}),
            (" may see; the filter decides what ", {}),
            ("each tool", {"italic": True, "colour": INK}),
            (" may see.\n", {}),
            ("And rest is not a leak.", {"bold": True, "colour": INK}),
            (" The NBA publishes its schedule in August. Only outcomes are gated.", {}),
        ],
        size=16,
        w=Inches(10.8),
    )

    # 6 — the tools -------------------------------------------------------
    s = new("The agent")
    headline(
        s,
        Inches(1.25),
        [
            ("Seven tools. That is its ", {}),
            ("entire", {"colour": ACCENT, "italic": True}),
            (" world.", {}),
        ],
        size=34,
    )
    widths = [Inches(4.3), Inches(5.0), Inches(2.1)]
    rows = [
        ("TOOL", "RETURNS", "STATE"),
        (
            ("retrieve_matchup_context", {"font": MONO}),
            "Ratings, rest, injuries, H2H as of a date",
            ("working", {"colour": GOOD, "bold": True}),
        ),
        (
            ("retrieve_team_form", {"font": MONO}),
            "Rolling 10-game record and point margin",
            ("working", {"colour": GOOD, "bold": True}),
        ),
        (
            ("retrieve_injuries", {"font": MONO}),
            "Who was known out that morning",
            ("working", {"colour": GOOD, "bold": True}),
        ),
        (
            ("retrieve_player_splits", {"font": MONO}),
            "Season averages, fatigue split",
            ("working", {"colour": GOOD, "bold": True}),
        ),
        (
            ("predict_win_probability", {"font": MONO}),
            "The model's number",
            ("working", {"colour": GOOD, "bold": True}),
        ),
        (
            ("retrieve_schedule", {"font": MONO}),
            "Forward slate",
            ("awaiting input", {"colour": SLATE}),
        ),
        (
            ("predict_stat_line", {"font": MONO}),
            "Projected pts / reb / ast",
            ("not built", {"colour": SLATE}),
        ),
    ]
    table(s, ML, Inches(2.15), CW, rows, widths, size=13)
    body(
        s,
        Inches(6.05),
        [
            ("No database. No web. ", {}),
            ("No way to invent a number.", {"bold": True, "colour": INK}),
            ("  The two unbuilt tools report their own gap rather than guessing.", {}),
        ],
        size=15,
    )

    # 7 — skills ----------------------------------------------------------
    s = new("Skills")
    headline(
        s, Inches(1.3), "The rules live in Markdown,\nnot buried in Python", size=34
    )
    body(
        s,
        Inches(2.85),
        "Each tool has a skill file — when to call it, how to read the answer. Loaded "
        "into the system prompt at startup, so a teammate changes agent behaviour "
        "without touching code.",
        size=16,
        w=Inches(10.2),
    )

    q = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ML, Inches(4.05), Inches(10.6), Inches(1.0)
    )
    q.fill.solid()
    q.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xF0)
    q.line.fill.background()
    q.shadow.inherit = False
    rule(s, Inches(4.05), w=Pt(4), thick=Inches(1.0))
    text(
        s,
        ML + Inches(0.45),
        Inches(4.32),
        Inches(9.8),
        Inches(0.5),
        "Do not apply your own injury penalty on top of the model's number.",
        size=18,
        colour=INK,
        italic=False,
        bold=True,
    )
    text(
        s,
        ML + Inches(0.45),
        Inches(4.72),
        Inches(9.8),
        Inches(0.3),
        "skills/retrieve_injuries.md",
        size=12,
        colour=MUTED,
        font=MONO,
    )

    body(
        s,
        Inches(5.45),
        [
            ("Why not a database? ", {"bold": True, "colour": INK}),
            (
                "We considered letting the agent query Postgres freely. Rejected for the "
                "same reason we removed the betting line: a free-form query surface is how "
                "an agent reaches data nobody intended it to see.",
                {},
            ),
        ],
        size=15,
        w=Inches(10.6),
    )

    # 8 — the model -------------------------------------------------------
    s = new("The model")
    headline(
        s,
        Inches(1.25),
        [
            ("Logistic regression — ", {}),
            ("chosen", {"colour": ACCENT, "italic": True}),
            (", not defaulted to", {}),
        ],
        size=32,
    )
    text(
        s,
        ML,
        Inches(2.15),
        Inches(6.2),
        Inches(0.4),
        "σ(w·x + b) over eight standardised features",
        size=18,
        colour=INK,
        font=MONO,
    )

    for i, (t, d) in enumerate(
        [
            ("Readable", "Weights can be argued with in a pull request"),
            ("Portable", "A few hundred bytes of named numbers, not a pickle"),
            (
                "Dependency-free",
                "Loads without sklearn — agent, UI and harness all run",
            ),
        ]
    ):
        yy = Inches(2.95) + i * Inches(0.78)
        rule(s, yy + Inches(0.09), w=Inches(0.22), thick=Pt(3))
        text(
            s,
            ML + Inches(0.42),
            yy,
            Inches(2.3),
            Inches(0.3),
            t,
            size=15,
            colour=INK,
            bold=True,
        )
        text(
            s,
            ML + Inches(2.55),
            yy + Inches(0.02),
            Inches(4.6),
            Inches(0.4),
            d,
            size=13,
            colour=MUTED,
        )

    text(
        s,
        Inches(8.2),
        Inches(2.5),
        Inches(4.2),
        Inches(0.3),
        "TOP FEATURE WEIGHTS",
        size=10,
        colour=MUTED,
        bold=True,
    )
    fw = [Inches(2.7), Inches(1.5)]
    frows = [
        (("win_pct_diff", {"font": MONO}), ("+0.396", {"bold": True})),
        (("form_margin_diff", {"font": MONO}), ("+0.378", {"bold": True})),
        (("injury_weight_diff", {"font": MONO}), ("−0.246", {"bold": True})),
        (("rest_diff", {"font": MONO}), ("+0.003", {"colour": SLATE})),
    ]
    table(s, Inches(8.2), Inches(2.9), Inches(4.2), frows, fw, header=False, size=13)

    # 9 — the trap --------------------------------------------------------
    s = new("The trap")
    headline(
        s,
        Inches(1.35),
        "The obvious way to build a feature\nis the leaking way",
        size=36,
    )
    body(
        s,
        Inches(3.0),
        [
            ("Season win percentage: group the season, take the mean. ", {}),
            (
                "That silently includes the game you are predicting.",
                {"bold": True, "colour": BAD},
            ),
            ("  The result is a very accurate and completely worthless model.", {}),
        ],
        size=17,
        w=Inches(10.6),
    )

    text(
        s,
        ML,
        Inches(4.35),
        Inches(8),
        Inches(0.3),
        "TESTS THAT CANNOT FAIL PROVE NOTHING — SO WE BROKE EACH RULE ON PURPOSE",
        size=10,
        colour=MUTED,
        bold=True,
    )
    trows = [
        (
            "Accumulators advanced early",
            ("3 tests caught it", {"bold": True, "colour": GOOD}),
        ),
        ("Form window drifted out of sync", ("1 test", {"bold": True, "colour": GOOD})),
        ("Test season added to training", ("3 tests", {"bold": True, "colour": GOOD})),
    ]
    table(
        s,
        ML,
        Inches(4.75),
        Inches(8.2),
        trows,
        [Inches(5.6), Inches(2.6)],
        header=False,
        size=14,
    )

    # 10 — the incident (inverted, full bleed) ----------------------------
    s = new(paper=INK)
    text(
        s, ML, MT, CW, Inches(0.3), "THE INCIDENT", size=11, colour=ACCENT_L, bold=True
    )
    headline(
        s,
        Inches(1.55),
        [("The agent cheated —\nand it wasn't even trying to", {"colour": PAPER})],
        size=40,
    )
    rule(s, Inches(3.55), w=Inches(1.1), colour=ACCENT_L)
    text(
        s,
        ML,
        Inches(4.05),
        Inches(11),
        Inches(0.9),
        "“The closing betting line favors the home team, ORL (-5.5 spread)”",
        size=24,
        colour=ACCENT_L,
        italic=True,
        spacing=1.25,
    )
    text(
        s,
        ML,
        Inches(5.0),
        Inches(10.5),
        Inches(0.3),
        "— the agent, quoting the exact benchmark we grade it against",
        size=13,
        colour=SLATE,
    )
    text(
        s,
        ML,
        Inches(5.75),
        Inches(11),
        Inches(0.8),
        [
            ("Telling a model not to peek is a request.  ", {"colour": SLATE}),
            ("Removing the tool is a guarantee.", {"colour": PAPER, "bold": True}),
        ],
        size=20,
        spacing=1.3,
    )

    # 11 — three arms -----------------------------------------------------
    s = new("The experiment")
    headline(
        s,
        Inches(1.3),
        [("Three arms differing by ", {}), ("exactly one tool", {"colour": ACCENT})],
        size=34,
    )
    arows = [
        ("ARM", "WHAT IT IS", "HAS THE MODEL?", "LLM?"),
        (
            ("A", {"bold": True, "size": 15}),
            "Model only",
            ("is the model", {"colour": MUTED}),
            "no",
        ),
        (
            ("B", {"bold": True, "size": 15}),
            "Agent alone, reasoning from tools",
            ("no", {"colour": BAD, "bold": True}),
            "yes",
        ),
        (
            ("C", {"bold": True, "size": 15}),
            "Agent + the model's number",
            ("yes", {"colour": GOOD, "bold": True}),
            "yes",
        ),
    ]
    table(
        s,
        ML,
        Inches(2.3),
        CW,
        arows,
        [Inches(1.1), Inches(5.6), Inches(2.9), Inches(1.8)],
        size=14,
        row_h=Inches(0.5),
    )
    body(
        s,
        Inches(4.5),
        [
            ("Same agent, same prompt, same data, same gate. ", {}),
            ("The difference is the measurement", {"bold": True, "colour": INK}),
            (
                " — and a test enforces that the tool lists differ by exactly one entry.",
                {},
            ),
        ],
        size=16,
        w=Inches(10.8),
    )
    text(
        s,
        ML,
        Inches(5.6),
        Inches(10),
        Inches(0.4),
        "Hypothesis, stated before the run:  C beats both.",
        size=18,
        colour=ACCENT,
        bold=True,
    )

    # 12 — results chart --------------------------------------------------
    s = new("Result")
    headline(
        s,
        Inches(1.2),
        [("Our hypothesis was ", {}), ("wrong", {"colour": BAD})],
        size=36,
    )
    text(
        s,
        ML,
        Inches(2.05),
        Inches(6.2),
        Inches(0.4),
        "Accuracy on all 1,322 games of 2025-26 — a season the model never trained on",
        size=14,
        colour=MUTED,
    )
    bar_chart(
        s,
        ML - Inches(0.2),
        Inches(2.5),
        Inches(6.6),
        Inches(4.0),
        ["Always home", "Our model", "Closing line"],
        [("accuracy %", (55.5, 66.5, 69.0))],
        [ACCENT],
        maximum=80,
    )
    text(
        s,
        Inches(7.5),
        Inches(2.6),
        Inches(4.9),
        Inches(0.4),
        "TRAIN 66.8%   ·   TEST 66.5%",
        size=12,
        colour=MUTED,
        bold=True,
    )
    text(
        s,
        Inches(7.5),
        Inches(3.05),
        Inches(4.9),
        Inches(0.6),
        "+0.3% generalisation gap. Not overfit.",
        size=16,
        colour=INK,
    )
    rule(s, Inches(3.95), x=Inches(7.5), w=Inches(0.9))
    text(
        s,
        Inches(7.5),
        Inches(4.3),
        Inches(4.9),
        Inches(1.8),
        [
            ("But arm C did not beat arm A.\n", {"colour": INK, "bold": True}),
            ("When the agent overruled the model it was wrong ", {}),
            ("15 times out of 19", {"colour": BAD, "bold": True}),
            (".\n", {}),
            (
                "two 40-game samples pooled · two-sided sign test p ≈ 0.019",
                {"size": 11, "colour": SLATE},
            ),
        ],
        size=15,
        colour=MUTED,
        spacing=1.35,
    )

    # 13 — before / after chart -------------------------------------------
    s = new("So we fixed it")
    headline(s, Inches(1.2), "Written rules recovered half the loss", size=34)
    text(
        s,
        ML,
        Inches(2.0),
        Inches(7.4),
        Inches(0.4),
        "Same games, same data. The only change was the prompt.",
        size=14,
        colour=MUTED,
    )
    bar_chart(
        s,
        ML - Inches(0.2),
        Inches(2.45),
        Inches(7.4),
        Inches(4.05),
        ["arm C accuracy", "overrides (of 80)"],
        [("before", (58.8, 19.0)), ("after", (66.3, 11.0))],
        [SLATE, ACCENT],
        maximum=70,
    )
    text(
        s,
        Inches(8.3),
        Inches(2.75),
        Inches(4.2),
        Inches(2.4),
        [
            ("Replicates on both seeds", {"colour": INK, "bold": True, "size": 17}),
            (", by the mechanism predicted: it overrules ", {}),
            ("less", {"italic": True, "colour": INK}),
            (".\n", {}),
            ("Log loss 0.674 → 0.591.", {}),
        ],
        size=15,
        colour=MUTED,
        spacing=1.35,
    )
    warn = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(4.7), Inches(4.2), Inches(1.35)
    )
    warn.fill.solid()
    warn.fill.fore_color.rgb = RGBColor(0xFE, 0xF6, 0xE7)
    warn.line.fill.background()
    warn.shadow.inherit = False
    text(
        s,
        Inches(8.55),
        Inches(4.92),
        Inches(3.8),
        Inches(1.0),
        [
            ("Not solved.\n", {"bold": True, "colour": INK}),
            (
                "The override test is now under-powered (p ≈ 0.23), not passed. "
                "Less overruling to measure ≠ fixed.",
                {},
            ),
        ],
        size=12,
        colour=MUTED,
        spacing=1.3,
    )

    # 14 — the money -------------------------------------------------------
    s = new("What it is worth")
    headline(
        s, Inches(1.2), [("Bet $100 a game. ", {}), ("Who ends up ahead?", {"colour": ACCENT})],
        size=34,
    )
    text(
        s, ML, Inches(2.05), Inches(7.6), Inches(0.4),
        "Full season, 1,322 games. Real spreads and scores; prices modelled with a 4.5% house margin.",
        size=13, colour=MUTED,
    )
    bar_chart(
        s, ML - Inches(0.2), Inches(2.5), Inches(7.2), Inches(4.0),
        ["Our predictor", "Always favourite", "Always home"],
        [("profit, $", (-3325, -5544, -8246))],
        [ACCENT],
    )
    text(
        s, Inches(8.4), Inches(2.6), Inches(4.1), Inches(0.4),
        "EVERYTHING LOSES", size=12, colour=BAD, bold=True,
    )
    text(
        s, Inches(8.4), Inches(3.05), Inches(4.1), Inches(1.0),
        "The house margin is the bar. None of our approaches clears it — and that is "
        "the honest answer.",
        size=15, colour=INK, spacing=1.3,
    )
    rule(s, Inches(4.35), x=Inches(8.4), w=Inches(0.9))
    text(
        s, Inches(8.4), Inches(4.7), Inches(4.1), Inches(2.0),
        [
            ("But look at the first two.\n", {"colour": INK, "bold": True}),
            ("Always backing the favourite wins ", {}),
            ("more games", {"colour": INK, "bold": True}),
            (" (69.0% vs 66.4%) and loses ", {}),
            ("more money", {"colour": BAD, "bold": True}),
            (".\n", {}),
            ("Favourites win often and pay badly.", {"size": 13}),
        ],
        size=14, colour=MUTED, spacing=1.35,
    )

    # 14 — the confound ---------------------------------------------------
    s = new("The measurement that changed our design")
    headline(s, Inches(1.3), "We tried to write “star out → drop N%”", size=34)
    body(
        s,
        Inches(2.15),
        "There is no N this data supports — and how that failed is the point.",
        size=16,
    )
    crows = [
        ("COMPARISON", "RESULT", "VERDICT"),
        (
            "All games, either star out",
            "home wins less when the AWAY star sits",
            ("backwards", {"colour": BAD, "bold": True}),
        ),
        (
            "Only teams that have a star",
            "+5.6% more wins without him (z = 2.6)",
            ("backwards, significant", {"colour": BAD, "bold": True}),
        ),
        (
            ("Each team vs itself", {"bold": True}),
            ("+0.0%  (se 3.3%, n = 21)", {"bold": True}),
            ("the real answer", {"colour": GOOD, "bold": True}),
        ),
    ]
    table(
        s,
        ML,
        Inches(2.9),
        CW,
        crows,
        [Inches(3.9), Inches(4.9), Inches(2.6)],
        size=13,
        row_h=Inches(0.52),
    )
    body(
        s,
        Inches(5.15),
        [
            ("Both flawed comparisons are confounded identically: ", {}),
            (
                "having a 20-ppg scorer is a property of good teams",
                {"bold": True, "colour": INK},
            ),
            (", and a good team is still good on the night its star sits.", {}),
        ],
        size=16,
        w=Inches(11),
    )
    text(
        s,
        ML,
        Inches(6.15),
        Inches(11),
        Inches(0.4),
        "So the rule says the opposite of what we set out to write: report the injury, let the model price it.",
        size=15,
        colour=ACCENT,
        bold=True,
    )

    # 15 — what we didn't do ----------------------------------------------
    s = new("Honest accounting")
    headline(s, Inches(1.3), "What we did not do", size=38)
    items = [
        (
            "predict_stat_line was never built",
            "a stated deliverable; the tool reports its own gap",
        ),
        ("n = 80 across two samples", "better than 40, still small"),
        (
            "Injury records are transaction dates",
            "not news timestamps — a same-day placement can appear",
        ),
        (
            "Star-to-team mapping is a season stale",
            "“injured” and “changed teams” are not cleanly separated",
        ),
        (
            "Two data sources have no upstream",
            "including the odds file the whole benchmark rests on",
        ),
        ("No opponent-adjusted strength", "probably the largest gain still available"),
    ]
    y = Inches(2.3)
    for t, d in items:
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL, ML, y + Inches(0.11), Inches(0.1), Inches(0.1)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        dot.shadow.inherit = False
        text(
            s,
            ML + Inches(0.35),
            y,
            Inches(5.4),
            Inches(0.35),
            t,
            size=15,
            colour=INK,
            bold=True,
        )
        text(
            s,
            ML + Inches(5.9),
            y + Inches(0.02),
            Inches(5.6),
            Inches(0.35),
            d,
            size=13,
            colour=MUTED,
        )
        y += Inches(0.68)

    # 16 — contributions --------------------------------------------------
    s = new("What we actually learned")
    headline(s, Inches(1.3), "Three things we did not expect", size=38)
    keeps = [
        (
            "The simple thing won",
            "A small statistical model beat both AI approaches. We assumed the opposite going in.",
        ),
        (
            "Adding AI made it worse",
            "And we could measure why: it talked itself out of correct answers, 15 times in 19.",
        ),
        (
            "Accuracy is not money",
            "Betting the favourite wins MORE games than our predictor and loses MORE money.",
        ),
    ]
    cw2 = Inches(3.6)
    for i, (t, d) in enumerate(keeps):
        cx = ML + i * (cw2 + Inches(0.42))
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(2.4), cw2, Inches(2.35)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF4, 0xF9, 0xF5)
        box.line.color.rgb = RGBColor(0xD5, 0xE8, 0xDA)
        box.shadow.inherit = False
        box.adjustments[0] = 0.04
        text(
            s,
            cx + Inches(0.34),
            Inches(2.75),
            cw2 - Inches(0.68),
            Inches(0.4),
            t,
            size=17,
            colour=GOOD,
            bold=True,
        )
        text(
            s,
            cx + Inches(0.34),
            Inches(3.3),
            cw2 - Inches(0.68),
            Inches(1.2),
            d,
            size=13,
            colour=INK,
            spacing=1.3,
        )

    body(
        s,
        Inches(5.3),
        [
            ("The scope of our claim is narrow: ", {}),
            ("this", {"italic": True, "colour": INK}),
            (" agent, with ", {}),
            ("these", {"italic": True, "colour": INK}),
            (" seven tools, on ", {}),
            ("these", {"italic": True, "colour": INK}),
            (
                " 80 paired games, degraded a good estimate. We have not shown that LLM "
                "agents cannot help.",
                {},
            ),
        ],
        size=16,
        w=Inches(11),
    )
    text(
        s,
        ML,
        Inches(6.35),
        Inches(11),
        Inches(0.4),
        "We set out to see what we could build. The useful part turned out to be measuring what didn't work.",
        size=16,
        colour=ACCENT,
        bold=True,
    )

    # 17 — live -----------------------------------------------------------
    s = new("Live")
    headline(s, Inches(1.3), "It runs from a clean clone", size=38)
    lrows = [
        (
            ("pytest", {"font": MONO}),
            "73 tests, incl. the mutation-tested leakage guards",
        ),
        (
            ("python -m models.train", {"font": MONO}),
            "Refits, prints the weights (~7s)",
        ),
        (
            ("python -m eval.three_arms", {"font": MONO}),
            "Arm A + both baselines, 1,322 games (~3s)",
        ),
        (
            ("python -m eval.injury_impact", {"font": MONO}),
            "The confounding story, reproduced",
        ),
        (
            ("python -m agent.run --model ollama", {"font": MONO}),
            "The real agent, local, no API key (~40s)",
        ),
        (
            ("python -m scripts.gate_snapshot", {"font": MONO}),
            "Materialises the gate, prints what it removed",
        ),
        (
            ("streamlit run ui/app.py", {"font": MONO}),
            "Report, tools, gating proof, build status",
        ),
    ]
    table(
        s, ML, Inches(2.3), CW, lrows, [Inches(5.5), Inches(5.9)], header=False, size=13
    )
    rule(s, Inches(6.15), w=Inches(1.1))
    text(
        s,
        ML,
        Inches(6.5),
        Inches(11),
        Inches(0.4),
        "github.com/joshcannonai/nba-game-intelligence-agent",
        size=16,
        colour=INK,
        font=MONO,
    )

    prs.save(str(path))
    return path


def main() -> None:
    out = (
        Path(sys.argv[1])
        if len(sys.argv) > 1
        else (Path.home() / "Desktop" / "cecs499-docs" / "CECS499-Presentation.pptx")
    )
    out.parent.mkdir(parents=True, exist_ok=True)
    print(build(out))


if __name__ == "__main__":
    main()
